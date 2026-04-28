#!/usr/bin/env python3
"""提取 raw/articles/ 中的二进制文件内容（PPTX/XLSX/VSDX）"""
import os, sys, zipfile, re
from pathlib import Path

BASE = Path(__file__).parent.parent / "raw" / "articles"
OUT  = Path(__file__).parent / "binary_extract"
OUT.mkdir(exist_ok=True)

files = {
    "pptx_mculess_audio":  BASE / "GPAN 车载MCULess和分布式音频介绍(1.8).pptx",
    "pptx_comm_08":        BASE / "GPAN 车载通信介绍-应用场景价值分析(0.8).pptx",
    "pptx_comm_06":        BASE / "GPAN 车载通信介绍-应用场景价值分析0.6.pptx",
    "pptx_seat":           BASE / "MCU-Less 座椅项目讨论_20260305-V0.5.pptx",
    "vsdx_scenarios":      BASE / "GPAN应用场景--外发理想.vsdx",
    "xlsx_cost":           BASE / "成本核算0.4 -- 公共.xlsx",
}

# ── PPTX 提取 ──────────────────────────────────────────────────────
def extract_pptx(path, key):
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            # 表格
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        slide_texts.append(" | ".join(cells))
        if slide_texts:
            lines.append(f"\n## Slide {i}")
            lines.extend(slide_texts)
    out_file = OUT / f"{key}.txt"
    out_file.write_text("\n".join(lines), encoding="utf-8")
    print(f"[PPTX] {path.name} → {len(prs.slides)} slides → {out_file.name}")

# ── XLSX 提取 ──────────────────────────────────────────────────────
def extract_xlsx(path, key):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"\n## Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append("\t".join(cells))
    out_file = OUT / f"{key}.txt"
    out_file.write_text("\n".join(lines), encoding="utf-8")
    print(f"[XLSX] {path.name} → {len(wb.worksheets)} sheets → {out_file.name}")

# ── VSDX 提取 ──────────────────────────────────────────────────────
def extract_vsdx(path, key):
    lines = []
    with zipfile.ZipFile(path, "r") as z:
        xml_files = [n for n in z.namelist() if n.endswith(".xml") or n.endswith(".rels")]
        lines.append(f"# VSDX Contents: {len(xml_files)} XML files")
        # 只读取 visio/pages/ 下的页面 XML
        page_files = [n for n in xml_files if "visio/pages/page" in n and n.endswith(".xml")]
        lines.append(f"# Page files: {page_files}")
        for pf in sorted(page_files):
            lines.append(f"\n## {pf}")
            content = z.read(pf).decode("utf-8", errors="replace")
            # 提取所有文本节点内容
            texts = re.findall(r"<t>(.*?)</t>", content)
            for t in texts:
                t = t.strip()
                if t:
                    lines.append(t)
    out_file = OUT / f"{key}.txt"
    out_file.write_text("\n".join(lines), encoding="utf-8")
    print(f"[VSDX] {path.name} → {len(page_files)} pages → {out_file.name}")

# ── 主流程 ──────────────────────────────────────────────────────────
for key, path in files.items():
    if not path.exists():
        print(f"[SKIP] 文件不存在: {path}")
        continue
    try:
        if path.suffix.lower() == ".pptx":
            extract_pptx(path, key)
        elif path.suffix.lower() == ".xlsx":
            extract_xlsx(path, key)
        elif path.suffix.lower() == ".vsdx":
            extract_vsdx(path, key)
    except Exception as e:
        print(f"[ERROR] {path.name}: {e}")

print("\n完成！输出目录:", OUT)
